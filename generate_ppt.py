"""
AutoPlayOSU! 專題進度報告 — 本次進度
只包含這次做的新改進
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# === 設計常數 ===
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_PINK = RGBColor(0xFF, 0x66, 0xAA)
ACCENT_BLUE = RGBColor(0x66, 0xBB, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xBB)
DARK_CARD = RGBColor(0x22, 0x22, 0x3A)
GREEN = RGBColor(0x66, 0xFF, 0x99)
ORANGE = RGBColor(0xFF, 0xAA, 0x44)
RED = RGBColor(0xFF, 0x66, 0x66)
YELLOW = RGBColor(0xFF, 0xDD, 0x55)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, font_name="Microsoft JhengHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return tf


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_accent_line(slide, top=1.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(2), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_PINK
    shape.line.fill.background()


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 1, 1.5, 11, 1, "AutoPlayOSU!", font_size=54, color=ACCENT_PINK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 2.7, 11, 0.8, "本次進度報告", font_size=30, color=WHITE, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.7), Inches(3.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

add_text(slide, 1, 4.2, 11, 0.6, "2026 年 2 月", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# 本次重點標籤
tags = ["OOM 修復", "分模型 Lookahead", "雙模型推理", "訓練問題排查"]
tag_start_x = 2.5
for i, tag in enumerate(tags):
    x = tag_start_x + i * 2.2
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.5), Inches(2), Inches(0.45))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CARD
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(1)
    card.adjustments[0] = 0.3
    card.text_frame.paragraphs[0].text = tag
    card.text_frame.paragraphs[0].font.size = Pt(12)
    card.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    card.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    card.text_frame.paragraphs[0].font.name = "Microsoft JhengHei"


# ============================================================
# Slide 2: 上次回顧 → 本次目標
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.8, "上次回顧 → 本次目標", font_size=36, color=ACCENT_PINK, bold=True)
add_accent_line(slide)

add_card(slide, 0.8, 1.6, 5.5, 3)
add_text(slide, 1.1, 1.7, 5, 0.5, "上次進度", font_size=22, color=GRAY, bold=True)
add_bullet_text(slide, 1.1, 2.3, 5, 2.2, [
    "✅ GPU 加速 (CPU → GPU 推理)",
    "✅ DXcam 高速截圖",
    "✅ TensorRT 模型量化 (FP16)",
    "✅ 移除 FPS 限制",
    "✅ 解析度匹配 (1080p)",
    "✅ 修復座標雙重歸一化 bug",
    "✅ 修復平均值回歸問題",
], font_size=15)

add_card(slide, 6.8, 1.6, 5.5, 3)
add_text(slide, 7.1, 1.7, 5, 0.5, "本次要解決", font_size=22, color=ACCENT_PINK, bold=True)
add_bullet_text(slide, 7.1, 2.3, 5, 2.2, [
    "🔴 訓練時記憶體不足 (OOM 崩潰)",
    "🔴 Actions 模型準確率只有 56%",
    "🔴 Combined 模型按鍵幾乎不作用",
    "🔴 雙模型推理 FPS 大幅下降",
    "🔴 訓練與實際遊戲表現有落差",
], font_size=15)

# 下方：架構圖提醒
add_card(slide, 0.8, 5.0, 11.5, 2.2)
add_text(slide, 1.1, 5.1, 11, 0.5, "系統流程提醒", font_size=18, color=ACCENT_BLUE, bold=True)

flow_items = [
    ("遊戲畫面", ACCENT_BLUE),
    ("螢幕截圖", ORANGE),
    ("影像處理", WHITE),
    ("AI 推理", GREEN),
    ("動作執行", ACCENT_PINK),
]
for i, (title, color) in enumerate(flow_items):
    x = 1.2 + i * 2.3
    add_card(slide, x, 5.7, 1.8, 0.7)
    add_text(slide, x + 0.1, 5.8, 1.6, 0.5, title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + 1.85, 5.85, 0.4, 0.4, "→", font_size=20, color=ACCENT_PINK, bold=True)


# ============================================================
# Slide 3: 修復 1 — 記憶體不足 (OOM)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.8, "修復 ① 記憶體不足 (OOM)", font_size=36, color=ACCENT_PINK, bold=True)
add_accent_line(slide)

# 問題
add_card(slide, 0.8, 1.6, 5.5, 2.5)
add_text(slide, 1.1, 1.7, 5, 0.5, "💥 問題", font_size=22, color=RED, bold=True)
add_bullet_text(slide, 1.1, 2.3, 5, 1.8, [
    "np.concatenate 把所有資料集合併成一個大陣列",
    "56,789 個 (10, 60, 80) 的影格",
    "→ 需要 5~10 GB 連續記憶體",
    "→ ArrayMemoryError 崩潰，無法訓練",
], font_size=15)

# 解法
add_card(slide, 6.8, 1.6, 5.5, 2.5)
add_text(slide, 7.1, 1.7, 5, 0.5, "✅ 解法", font_size=22, color=GREEN, bold=True)
add_bullet_text(slide, 7.1, 2.3, 5, 1.8, [
    "OsuLazyDataset — 懶加載架構",
    "只存索引，__getitem__ 時才讀取",
    "完全不合併大陣列（零額外記憶體）",
    "即時轉 float16 (192KB → 96KB/幀)",
], font_size=15)

# 技術細節
add_card(slide, 0.8, 4.5, 7, 2.8)
add_text(slide, 1.1, 4.6, 6.5, 0.5, "技術細節", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_text(slide, 1.1, 5.2, 6.5, 2, [
    "1. 分塊存儲：每個 dataset 保持獨立的 numpy 陣列（chunk）",
    "2. 累積索引：O(1) 查找任意全局索引對應的 chunk + 本地索引",
    "3. 增量陣列建立：np.empty 預分配 + 逐幀填入 + 逐幀釋放",
    "4. 索引過採樣：RandomOverSampler 只複製索引，不複製圖像",
], font_size=14)

# 結果
add_card(slide, 8.2, 4.5, 4.1, 2.8)
add_text(slide, 8.5, 4.6, 3.5, 0.5, "📊 結果", font_size=20, color=GREEN, bold=True)

add_text(slide, 8.5, 5.3, 3.5, 0.4, "記憶體使用", font_size=14, color=GRAY)
add_text(slide, 8.5, 5.7, 1.5, 0.5, "10 GB", font_size=28, color=RED, bold=True)
add_text(slide, 10, 5.85, 0.5, 0.3, "→", font_size=22, color=ACCENT_PINK, bold=True)
add_text(slide, 10.4, 5.7, 1.5, 0.5, "2 GB", font_size=28, color=GREEN, bold=True)

add_text(slide, 8.5, 6.4, 3.5, 0.4, "✅ 訓練不再崩潰", font_size=16, color=WHITE, bold=True)


# ============================================================
# Slide 4: 修復 2 — 分模型 Lookahead
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.8, "修復 ② 分模型 Lookahead", font_size=36, color=ACCENT_PINK, bold=True)
add_accent_line(slide)

# 問題
add_card(slide, 0.8, 1.6, 5.5, 2.3)
add_text(slide, 1.1, 1.7, 5, 0.5, "💥 問題", font_size=22, color=RED, bold=True)
add_bullet_text(slide, 1.1, 2.3, 5, 1.5, [
    "所有模型都用 lookahead=6 (預測 200ms 後)",
    "Aim 沒問題（滑鼠移動連續可預測）",
    "但 Actions 模型只有 56% 準確率！",
    "按鍵是瞬間事件，200ms 太遠了",
], font_size=15)

# 解法
add_card(slide, 6.8, 1.6, 5.5, 2.3)
add_text(slide, 7.1, 1.7, 5, 0.5, "✅ 解法：因地制宜", font_size=22, color=GREEN, bold=True)
add_text(slide, 7.1, 2.3, 5, 0.4, "每種模型使用最適合的 lookahead：", font_size=15, color=WHITE)

# 表格
headers_x = [7.3, 9.2, 10.5]
add_text(slide, 7.3, 2.75, 1.5, 0.3, "模型", font_size=15, color=GRAY, bold=True)
add_text(slide, 9.2, 2.75, 1.5, 0.3, "Lookahead", font_size=15, color=GRAY, bold=True)
add_text(slide, 10.5, 2.75, 1.5, 0.3, "原因", font_size=15, color=GRAY, bold=True)

rows = [
    ("Aim", "6 幀 (~200ms)", "連續移動", ACCENT_BLUE),
    ("Actions", "2 幀 (~67ms)", "瞬間事件", GREEN),
    ("Combined", "3 幀 (~100ms)", "折衷", ORANGE),
]
for i, (m, la, reason, color) in enumerate(rows):
    y = 3.1 + i * 0.3
    add_text(slide, 7.3, y, 1.5, 0.3, m, font_size=14, color=color, bold=True)
    add_text(slide, 9.2, y, 1.5, 0.3, la, font_size=14, color=WHITE)
    add_text(slide, 10.5, y, 1.5, 0.3, reason, font_size=13, color=GRAY)

# 說明圖
add_card(slide, 0.8, 4.3, 11.5, 3)
add_text(slide, 1.1, 4.4, 11, 0.5, "Lookahead 原理", font_size=20, color=ACCENT_BLUE, bold=True)

add_text(slide, 1.1, 5.0, 11, 0.4, "把圖像和標籤錯開，讓模型學習「預測未來」以補償推理延遲", font_size=15, color=WHITE)

# 時間線示意
add_text(slide, 1.3, 5.6, 1.5, 0.3, "圖像 (輸入)", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, 1.3, 5.95, 1.5, 0.3, "標籤 (目標)", font_size=14, color=GREEN, bold=True)

for i in range(8):
    x = 3 + i * 1
    frame_color = ACCENT_BLUE if i < 6 else RGBColor(0x33, 0x33, 0x55)
    label_color = GREEN if i >= 2 else RGBColor(0x33, 0x33, 0x55)
    
    fc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.55), Inches(0.8), Inches(0.35))
    fc.fill.solid()
    fc.fill.fore_color.rgb = DARK_CARD
    fc.line.color.rgb = frame_color
    fc.line.width = Pt(1.5)
    fc.adjustments[0] = 0.2
    fc.text_frame.paragraphs[0].text = f"t{i}"
    fc.text_frame.paragraphs[0].font.size = Pt(11)
    fc.text_frame.paragraphs[0].font.color.rgb = frame_color
    fc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    lc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.95), Inches(0.8), Inches(0.35))
    lc.fill.solid()
    lc.fill.fore_color.rgb = DARK_CARD
    lc.line.color.rgb = label_color
    lc.line.width = Pt(1.5)
    lc.adjustments[0] = 0.2
    lc.text_frame.paragraphs[0].text = f"t{i+2}" if i < 6 else ""
    lc.text_frame.paragraphs[0].font.size = Pt(11)
    lc.text_frame.paragraphs[0].font.color.rgb = label_color
    lc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text(slide, 3, 6.5, 8, 0.4, "↑ lookahead=2 時：用 t0 的圖像學習 t2 的操作（預測 67ms 後）", font_size=13, color=GRAY)

# 快取不覆蓋
add_text(slide, 1.1, 6.85, 11, 0.4, "💡 快取檔名加入 lookahead 值（la2-dataset1.npz vs la6-dataset1.npz），不同設定互不覆蓋", font_size=13, color=YELLOW)


# ============================================================
# Slide 5: 修復 3 — DualEvalThread
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.8, "修復 ③ 雙模型推理優化", font_size=36, color=ACCENT_PINK, bold=True)
add_accent_line(slide)

# 問題
add_card(slide, 0.8, 1.6, 5.5, 2.2)
add_text(slide, 1.1, 1.7, 5, 0.5, "💥 問題", font_size=22, color=RED, bold=True)
add_text(slide, 1.1, 2.2, 5, 0.4, "同時跑 Aim + Actions 兩個模型時", font_size=15, color=WHITE)
add_bullet_text(slide, 1.1, 2.6, 5, 1, [
    "兩個線程各自呼叫 mss.grab() 截圖",
    "搶同一個 DWM 資源 → 互相拖慢",
    "FPS 從 29 暴跌到 19",
], font_size=15, color=ORANGE)

# 解法
add_card(slide, 6.8, 1.6, 5.5, 2.2)
add_text(slide, 7.1, 1.7, 5, 0.5, "✅ 解法：DualEvalThread", font_size=22, color=GREEN, bold=True)
add_text(slide, 7.1, 2.2, 5, 0.4, "一個線程統一截圖，序列跑兩個模型", font_size=15, color=WHITE)
add_bullet_text(slide, 7.1, 2.6, 5, 1, [
    "一次 mss.grab() 共用同一張截圖",
    "Aim 推理 → Actions 推理（序列執行）",
    "FPS 提升到 21",
], font_size=15, color=GREEN)

# 對比圖
add_card(slide, 0.8, 4.2, 5.5, 3.1)
add_text(slide, 1.1, 4.3, 5, 0.5, "Before：兩線程搶資源", font_size=18, color=RED, bold=True)

steps_old = [
    ("AimThread", "mss.grab() 30ms", RED),
    ("AimThread", "Aim推理 4ms", ACCENT_BLUE),
    ("ActionsThread", "mss.grab() 30ms", RED),
    ("ActionsThread", "Actions推理 4ms", GREEN),
]
for i, (thread, step, color) in enumerate(steps_old):
    y = 4.95 + i * 0.5
    add_text(slide, 1.3, y, 2, 0.3, thread, font_size=12, color=GRAY, font_name="Consolas")
    bar_width = 3 if "30ms" in step else 0.5
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), Inches(y), Inches(bar_width), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.3
    bar.text_frame.paragraphs[0].text = step
    bar.text_frame.paragraphs[0].font.size = Pt(10)
    bar.text_frame.paragraphs[0].font.color.rgb = WHITE
    bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text(slide, 1.3, 6.9, 5, 0.3, "兩次截圖 = 60ms 浪費 → 19 FPS", font_size=14, color=RED, bold=True)

add_card(slide, 6.8, 4.2, 5.5, 3.1)
add_text(slide, 7.1, 4.3, 5, 0.5, "After：共用截圖", font_size=18, color=GREEN, bold=True)

steps_new = [
    ("DualThread", "mss.grab() 30ms", ORANGE),
    ("DualThread", "Aim推理 4ms", ACCENT_BLUE),
    ("DualThread", "Actions推理 4ms", GREEN),
    ("DualThread", "numpy處理 8ms", GRAY),
]
for i, (thread, step, color) in enumerate(steps_new):
    y = 4.95 + i * 0.5
    add_text(slide, 7.3, y, 2, 0.3, thread, font_size=12, color=GRAY, font_name="Consolas")
    bar_width = 3 if "30ms" in step else (0.8 if "8ms" in step else 0.5)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(y), Inches(bar_width), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.3
    bar.text_frame.paragraphs[0].text = step
    bar.text_frame.paragraphs[0].font.size = Pt(10)
    bar.text_frame.paragraphs[0].font.color.rgb = WHITE
    bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text(slide, 7.3, 6.9, 5, 0.3, "一次截圖 = 46ms → 21 FPS ✅", font_size=14, color=GREEN, bold=True)


# ============================================================
# Slide 6: 訓練結果與問題分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.8, "訓練結果與問題分析", font_size=36, color=ACCENT_PINK, bold=True)
add_accent_line(slide)

# 訓練結果表格
add_card(slide, 0.8, 1.6, 11.5, 2.8)
add_text(slide, 1.1, 1.7, 11, 0.5, "訓練結果", font_size=20, color=ACCENT_BLUE, bold=True)

headers = ["模型", "訓練 Acc", "驗證 Acc", "Epochs", "狀態"]
col_xs = [1.3, 3.8, 5.8, 7.8, 9.3]
for i, h in enumerate(headers):
    add_text(slide, col_xs[i], 2.25, 2, 0.3, h, font_size=15, color=GRAY, bold=True)

sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(2.6), Inches(10.5), Inches(0.02))
sep.fill.solid()
sep.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x55)
sep.line.fill.background()

results = [
    ("Aim (la=6)", "~99%", "~70%", "~30", "✅ 可用", GREEN),
    ("Actions (la=6, 舊)", "99.5%", "56.8%", "32", "❌ 過擬合", RED),
    ("Actions (la=2, 新)", "~95%", "~65%", "~30", "⬆️ 改善中", YELLOW),
    ("Combined (la=3)", "~90%", "~60%", "~30", "⚠️ 按鍵弱", ORANGE),
]
for i, (model, train, val, ep, status, color) in enumerate(results):
    y = 2.75 + i * 0.38
    vals = [model, train, val, ep, status]
    for j, v in enumerate(vals):
        c = color if j == 0 or j == 4 else WHITE
        add_text(slide, col_xs[j], y, 2, 0.35, v, font_size=15, color=c, bold=(j==0))

# 核心問題
add_card(slide, 0.8, 4.8, 5.5, 2.5)
add_text(slide, 1.1, 4.9, 5, 0.5, "🔍 發現的核心問題", font_size=20, color=RED, bold=True)
add_bullet_text(slide, 1.1, 5.5, 5, 1.8, [
    "訓練準確度有 60-70%",
    "但實際上機遊戲表現很差",
    "看不出來有這個準確度",
    "",
    "→ 正在找原因中",
], font_size=16, color=YELLOW)

# 懷疑原因
add_card(slide, 6.8, 4.8, 5.5, 2.5)
add_text(slide, 7.1, 4.9, 5, 0.5, "📋 懷疑的原因", font_size=20, color=ORANGE, bold=True)
add_bullet_text(slide, 7.1, 5.5, 5, 1.8, [
    "1. Lookahead 還需微調最佳值",
    "2. Combined 的 MSE Loss 不適合分類",
    "3. 過採樣造成訓練 bias",
    "4. 驗證集 Acc 不等於實際遊戲 Acc",
    "5. 資料量/多樣性不足",
], font_size=14)


# ============================================================
# Slide 7: 下一步
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.8, "下一步計畫", font_size=36, color=ACCENT_PINK, bold=True)
add_accent_line(slide)

plans = [
    ("模型改進", [
        "• Combined 改用 BCE Loss",
        "  處理按鍵分類任務",
        "• 嘗試更大骨幹",
        "  (EfficientNet-B2)",
        "• 加入 Attention 機制",
        "  聚焦圓圈區域",
    ], ACCENT_BLUE),
    ("調查表現落差", [
        "• 分析推理時的實際",
        "  輸出 vs 訓練時預測",
        "• 錄製推理過程",
        "  與真實操作對比",
        "• 驗證 lookahead",
        "  的最佳數值",
    ], ORANGE),
    ("資料改進", [
        "• 收集更多不同星級",
        "  的訓練資料",
        "• 確保資料品質",
        "  (避免錯誤標註)",
        "• 嘗試 Data",
        "  Augmentation",
    ], GREEN),
]

for i, (title, items, color) in enumerate(plans):
    x = 0.8 + i * 4
    add_card(slide, x, 1.6, 3.6, 4)
    add_text(slide, x + 0.2, 1.7, 3.2, 0.5, title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    sep_s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.3), Inches(2.25), Inches(3), Inches(0.02))
    sep_s.fill.solid()
    sep_s.fill.fore_color.rgb = color
    sep_s.line.fill.background()
    
    add_bullet_text(slide, x + 0.3, 2.4, 3, 3, items, font_size=14)

# 本次進度總結
add_card(slide, 0.8, 6.0, 11.5, 1.2)
add_text(slide, 1.1, 6.1, 11, 0.4, "本次進度總結", font_size=18, color=ACCENT_PINK, bold=True)
summary_items = [
    ("✅ OOM 修復", "10GB → 2GB", GREEN),
    ("✅ 分模型 Lookahead", "Actions 56→65%", GREEN),
    ("✅ DualEvalThread", "19→21 FPS", GREEN),
    ("🔄 表現落差調查", "進行中", YELLOW),
]
for i, (item, detail, color) in enumerate(summary_items):
    x = 1.2 + i * 2.8
    add_text(slide, x, 6.5, 2, 0.25, item, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, 6.75, 2, 0.25, detail, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 8: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 1, 2.5, 11, 1, "Thank You", font_size=56, color=ACCENT_PINK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.8, "AutoPlayOSU!", font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.8), Inches(3.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

add_text(slide, 1, 5.5, 11, 0.5, "Q & A", font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# === 保存 ===
output_path = r"D:\Tools\gemini-cli-main\AutoPlayOSU!\AutoPlayOSU_Report.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
